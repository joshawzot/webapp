# route_handlers.py

from run import app, cache, redis_client
from db_operations import *
from tools_for_plots import get_full_table_data, plot_individual_points_map  # Add plot_individual_points_map to the import
from flask_caching import Cache

# Standard library imports
import os, base64, json, time
from io import BytesIO

# External libraries
import pandas as pd
import mysql.connector
from flask import Flask, request, make_response, redirect, url_for, session, send_file, render_template, render_template_string, jsonify
from pptx import Presentation
import zipfile
import numpy as np
import csv
import io
from PIL import Image
from sqlalchemy import create_engine
import traceback
import subprocess
import sys  # Add this import
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import threading

import re
from flask import render_template_string
import pandas as pd
from scipy.io import loadmat
import h5py
import numpy as np
from io import BytesIO
import scipy.io
# Create a lock for thread-safe plotting
#plot_lock = threading.Lock()

# Custom module imports
from generate_plot import generate_plot
#from generate_plot_ber_by_bls import generate_plot_ber_by_bls
from generate_plot_read_stability import generate_plot_read_stability

#from flask_caching import Cache
#cache = Cache(app, config={'CACHE_TYPE': 'simple'})

@app.route('/')
def home():
    username = session.get('username')
    print(username)
    if username:
        try:
            conn = create_connection()
            cursor = conn.cursor()
            databases = get_all_databases(cursor)
            
            # Try to get cached size first
            #total_size = cache.get('total_db_size')
            total_size = None
            if total_size is None:
                # If not in cache, calculate and store it
                _, total_size = get_total_database_size()
                cache.set('total_db_size', total_size, timeout=3600)  # Cache for 1 hr
                
            cursor.close()
            conn.close()
            return render_template('home_page.html', databases=databases, username=username, total_size=total_size)
        except mysql.connector.Error as err:
            return str(err), 500
    else:
        return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        if username:
            session['username'] = username
            return redirect(url_for('home'))
        return render_template('login.html', error="Please enter a username")
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.pop('username', None)
    return redirect(url_for('login'))

@app.route('/create-db')
def create_db_page():
    conn = create_connection()
    cursor = conn.cursor()
    databases = get_all_databases(cursor)  # Fetch all database names
    cursor.close()
    conn.close()

    return render_template('create_db_page.html')

@app.route('/save-txt-content/<database>/<table_name>', methods=['POST'])
def save_txt_content(database, table_name):
    try:
        content = request.json['content']
        connection = create_connection(database)
        cursor = connection.cursor()
        query = f"UPDATE `{table_name}` SET content = %s WHERE content IS NOT NULL LIMIT 1"
        cursor.execute(query, (content,))
        connection.commit()
        close_connection()
        return "Content saved successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

@app.route('/list-tables', methods=['POST', 'GET'])
def list_tables():
    if request.method == 'POST':
        session['database'] = request.form.get('database')

    database = session.get('database')

    tables = fetch_tables(database)  # Retrieve table data from the database
    table_names = ','.join(table['table_name'] for table in tables)
    print("table_names:", table_names)

    plot_function = "None"  # This could also be dynamically set based on POST or other conditions

    return render_template('list_tables.html', tables=tables, table_names=table_names, database=database, plot_function=plot_function)

@app.route('/view-table/<database>/<table_name>', methods=['GET'])
def view_table(database, table_name):
    """View the content of a specific table."""
    print('database:', database)
    print('table_name:', table_name)
    try:
        connection = create_connection(database)
        cursor = connection.cursor()
        if table_name.endswith('_txt'):
            query = f"SELECT content FROM `{table_name}` LIMIT 1"
            results = fetch_data(cursor, query)
            close_connection()
            content = results[0][0] if results else ''
            print('------')
            print(table_name)
            return render_template('table_txt.html', database=database, content=content, table_name=table_name)
        else:
            query = f"SELECT * FROM `{table_name}`"
            results = fetch_data(cursor, query)
            column_names = [desc[0] for desc in cursor.description]
            close_connection()
            return render_template('table.html', results=results, column_names=column_names)
    except mysql.connector.Error as err:
        return str(err)

# Define a dictionary to map plot function names to their corresponding functions
generate_plot_functions = {
    "generate_plot": generate_plot,
    #"generate_plot_ber_by_bls": generate_plot_ber_by_bls,
    "generate_plot_read_stability": generate_plot_read_stability,
}

@app.route('/render-plot/<database>/<table_name>/<plot_function>')
def render_plot(database, table_name, plot_function):
    try:
        if 'username' not in session:
            return "User not logged in", 403

        # Parse form data
        form_data_json = request.args.get('form_data', '{}')
        try:
            form_data = json.loads(form_data_json)
        except json.JSONDecodeError:
            return "Error: Invalid form data", 400

        # Map plot function names to actual functions
        plot_functions = {
            'generate_plot': generate_plot,
            'generate_plot_read_stability': generate_plot_read_stability,
            #'generate_plot_ber_by_bls': generate_plot_ber_by_bls,
        }

        # Validate plot function
        plot_function_impl = plot_functions.get(plot_function)
        if plot_function_impl is None:
            return f"Invalid plot function: {plot_function}", 400

        # Generate plot data with thread isolation
        try:
            # Create a new figure manager for this request
            import matplotlib.pyplot as plt
            plt.switch_backend('Agg')
            
            if plot_function == 'generate_plot':
                (plot_data,
                 sorted_table_names,
                 sorted_table_names_100ppm,
                 sorted_table_names_200ppm,
                 sorted_table_names_500ppm,
                 sorted_table_names_1000ppm,
                 best_32,
                 best_32_with_io,
                 outlier_coordinates,
                 correlation_analysis,
                 cluster_map,
                 sigma_distances,
                 num_states,
                 table_names,
                 sigma_table,
                 sigma_points) = plot_function_impl(table_name.split(','), database, form_data)
            else:
                plot_data = plot_function_impl(table_name.split(','), database, form_data)
                sorted_table_names = sorted_table_names_100ppm = \
                    sorted_table_names_200ppm = sorted_table_names_500ppm = \
                    sorted_table_names_1000ppm = best_32 = best_32_with_io = \
                    outlier_coordinates = correlation_analysis = cluster_map = \
                    sigma_distances = num_states = table_names = sigma_table = sigma_points = None

            # Clean up all figures created during this request
            plt.close('all')

            if plot_data is None:
                return "Failed to generate plot data", 400

            # Get the dimensions from the first table's data matrix
            first_table_name = table_name.split(',')[0]
            data_matrix, data_matrix_size = get_full_table_data(first_table_name, database)
            rows, cols = data_matrix_size
            
            # Generate individual points map
            print("Generating points map...")
            print("Outlier coordinates:", outlier_coordinates)
            points_map = plot_individual_points_map(outlier_coordinates, (rows, cols))
            print("Points map generated:", points_map is not None)

            # Debug print for template variables
            print("Template variables:")
            print("- cluster_map present:", points_map is not None)
            print("- outlier_coordinates present:", bool(outlier_coordinates))
            print("- outlier_coordinates length:", len(outlier_coordinates) if outlier_coordinates else 0)

            return render_template(
                'plot.html',
                plot_data=plot_data,
                sorted_table_names=sorted_table_names,
                sorted_table_names_100ppm=sorted_table_names_100ppm,
                sorted_table_names_200ppm=sorted_table_names_200ppm,
                sorted_table_names_500ppm=sorted_table_names_500ppm,
                sorted_table_names_1000ppm=sorted_table_names_1000ppm,
                best_32=best_32,
                best_32_with_io=best_32_with_io,
                outlier_coordinates=outlier_coordinates,
                correlation_analysis=correlation_analysis,
                cluster_map=points_map,
                sigma_distances=sigma_distances,
                num_states=num_states,
                table_names=table_names,
                target_values=form_data.get('target_values', []),
                sigma_table=sigma_table,
                sigma_points=sigma_points
            )

        except Exception as e:
            print(f"Error generating plot: {e}")
            return f"Error generating plot: {str(e)}", 500

    except Exception as e:
        return f"Error: {str(e)}", 500

@app.route('/download_csv/<unique_id>/<data_type>')
def download_csv2(unique_id, data_type):
    # Retrieve plot data either from cache or Redis
    cache_key = f"plot_data_{unique_id}"
    plot_data = cache.get(cache_key)
    if not plot_data:
        stored_data_json = redis_client.get(unique_id)
        if not stored_data_json:
            return "Error: Data not found", 404
        stored_data = json.loads(stored_data_json)
        database = stored_data["database"]
        table_names = stored_data["table_name"].split(',')
        form_data = stored_data["form_data"]
        plot_function = stored_data["plot_function"]
        generate_plot_function = generate_plot_functions.get(plot_function)
        if not generate_plot_function:
            return "Error: Invalid plot function selection", 400
        plot_data = generate_plot_function(table_names, database, form_data)

    # Generate CSV based on plot_data and data_type
    if data_type == "avg_std":
        avg_values, std_values, table_names, selected_groups = plot_data
        header = ["State"] + [f"{table_name}" for table_name in table_names] + ["Row Avg", "Row Std Dev"]
        table_data = [header]
        column_data = [[] for _ in table_names]

        for i, group in enumerate(selected_groups):
            row = [f"State {group}"]
            row_data = []

            for j, table_avg in enumerate(avg_values):
                avg = table_avg[i]
                row.append(f"{avg:.2f}")
                row_data.append(avg)
                column_data[j].append(avg)

            row_avg = np.mean(row_data)
            row_std = np.std(row_data)
            row.extend([f"{row_avg:.2f}", f"{row_std:.2f}"])
            table_data.append(row)

        col_avgs = [np.mean(col) for col in column_data]
        col_stds = [np.std(col) for col in column_data]
        table_data.append(["Col Avg"] + [f"{avg:.2f}" for avg in col_avgs] + ["-", "-"])
        table_data.append(["Col Std Dev"] + [f"{std:.2f}" for std in col_stds] + ["-", "-"])
        return generate_csv_response(table_data, "avg_std_data.csv")

    elif data_type in ["sigma", "ppm", "us"]:
        ber_results, _ = plot_data
        headers = ["State/Transition"] + [name for name in table_names] + ["Row Avg"]
        data_collections = [headers[:], headers[:], headers[:]]

        grouped_data = {}
        for entry in ber_results:
            key = entry[1]
            if key not in grouped_data:
                grouped_data[key] = []
            grouped_data[key].append((entry[2], entry[3], entry[4]))

        for key, values in grouped_data.items():
            rows = [[key], [key], [key]]
            for val in values:
                rows[0].append(f"{val[0]:.4f}")
                rows[1].append(f"{int(val[1])}")
                rows[2].append(f"{int(val[2])}")
            for row in rows:
                avg = np.mean([float(v) for v in row[1:]])
                row.append(f"{avg:.4f}")

            data_collections[0].append(rows[0])
            data_collections[1].append(rows[1])
            data_collections[2].append(rows[2])

        index = {"sigma": 0, "ppm": 1, "us": 2}[data_type]
        filename = f"{data_type}_data.csv"
        return generate_csv_response(data_collections[index], filename)

def generate_csv_response(data, filename):
    csv_output = StringIO()
    for row in data:
        csv_output.write(','.join(str(item) for item in row) + '\n')
    csv_output.seek(0)
    response = make_response(csv_output.getvalue())
    response.headers["Content-Disposition"] = f"attachment; filename={filename}"
    response.headers["Content-type"] = "text/csv"
    return response

@app.route('/download_csv')
def download_csv():
    database = request.args.get('database')
    table_name = request.args.get('table_name')

    try:
        # Generate the CSV data
        csv_data = get_csv_from_table(database, table_name)
        if csv_data is None:
            return "Error generating CSV file", 500

        # Create a response with the CSV data as a downloadable file
        response = make_response(csv_data)
        response.headers['Content-Disposition'] = f'attachment; filename={table_name}.csv'
        response.mimetype = 'text/csv'
        return response
    except Exception as e:
        return str(e), 500

@app.route('/download_npy')
def download_npy():
    database = request.args.get('database')
    table_name = request.args.get('table_name')
    try:
        data = get_npy_from_table(database, table_name)
        if data is None:
            return "Error retrieving data or no data available", 500

        bio = io.BytesIO(data)

        return send_file(
            bio,
            as_attachment=True,
            download_name=f'{table_name}.npy',
            mimetype='application/octet-stream'
        )
    except Exception as e:
        print(f"Download error: {e}")
        return str(e), 500

@app.route('/view-plot/<database>/<table_name>/<plot_function>', methods=['GET', 'POST'])
def view_plot(database, table_name, plot_function):
    print("view_plot")
    if request.method == "POST":
        print("POST:::::::::::::::::::::::::::::::::")
        plot_function_choice = request.form.get('plot_choice')
        if plot_function_choice:
            plot_function = plot_function_choice
            if plot_function in generate_plot_functions:
                if plot_function == "generate_plot":
                    return render_template('input_form_generate_plot.html', database=database, table_name=table_name, plot_function=plot_function)
                elif plot_function == "generate_plot_read_stability":
                    return render_template('input_form_generate_plot_read_stability.html', database=database, table_name=table_name, plot_function=plot_function)
                #elif plot_function == "generate_plot_ber_by_bls":
                    #return render_template('input_form_generate_plot_ber_by_bls.html', database=database, table_name=table_name, plot_function=plot_function)
            else:
                return f"Invalid plot function selection", 400

        if plot_function:
            print(f"plot_function: {plot_function}")
            # if plot_function in generate_plot_functions:
            #     form_data = get_form_data_generate_plot(request.form)
            #     form_data_json = json.dumps(form_data)

            if plot_function in ["generate_plot", "generate_plot_read_stability", '''"generate_plot_ber_by_bls"''']:
                # Determine the appropriate function to call based on plot_function
                if plot_function == "generate_plot":
                    form_data = get_form_data_generate_plot(request.form)
                elif plot_function == "generate_plot_read_stability":
                    form_data = get_form_data_generate_plot_read_stability(request.form)
                #elif plot_function == "generate_plot_ber_by_bls":
                    #form_data = get_form_data_generate_plot_ber_by_bls(request.form)
        
                # Convert form data to JSON
                form_data_json = json.dumps(form_data)

                # Redirect with the form data in the query string
                return redirect(f"/render-plot/{database}/{table_name}/{plot_function}?form_data={form_data_json}")
            else:
                return f"Invalid plot function selection", 400
        else:
            return f"Plot function not selected", 400
    else:
        table_names = table_name.split(',')
        print(table_names)
        print("GET:::::::::::::::::::::::::::::::::")
        return render_template('choose_plot_function_form.html')

from sqlalchemy import text
import pandas as pd

@app.route('/upload-file', methods=['POST'])
def upload_file():   #auto upload
    print("upload_file()")
    print("request.form:", request.form)
    print("request.files:", request.files)

    if 'db_name' not in request.form:
        return "No database selected", 400

    if 'files[]' not in request.files:
        print("No files part in request.files")  # Log missing files part
        return "No files part in the request", 400

    files = request.files.getlist('files[]')
    print("All files:", files)  # Log all files before filtering

    files = [f for f in files if f.filename]
    if not files:
        print("No files detected")  # Log no files detected
        return "No files detected", 400

    for file in files:
        print(f"File: {file.filename}, MIME Type: {file.mimetype}")

    db_name = request.form['db_name']
    engine = create_db_engine(db_name)

    results = []
    for file in files:
        filename = sanitize_table_name(file.filename)
        file_extension = filename.rpartition('_')[-1]
        print("file_extension:", file_extension)
        file_stream = BytesIO(file.read())

        try:
            df = process_file(file_stream, file_extension, db_name)
            print("Original DataFrame shape (rows, columns):", df.shape) #1296,1024

            # Final check before uploading
            if df.isnull().values.any():
                print("DataFrame contains NaN values before uploading.")
            else:
                print("DataFrame does not contain NaN values before uploading.")

            if not df.empty:
                #df = df.head(500)
                #df = df.iloc[:, :500] #1296,500

                print("DataFrame shape (rows, columns):", df.shape)
                df.to_sql(filename, engine, if_exists='replace', index=False)

                results.append(f"{filename} uploaded successfully")
            else:
                results.append(f"No data to upload for {filename}. Dataframe is empty.")
        except Exception as e:
            error_msg = f"Error processing {filename}: {str(e)}"
            results.append(error_msg)

    return jsonify(results=results)

@app.route('/delete-record/<database>/<table_name>', methods=['DELETE'])  # delete a table
def delete_record(database, table_name):
    try:
        connection = create_connection(database)
        cursor = connection.cursor()

        query = f"DROP TABLE `{table_name}`"
        cursor.execute(query)

        connection.commit()
        close_connection()

        return "Record deleted successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

@app.route('/delete-records/<database>', methods=['DELETE'])  # delete multiple tables
def delete_records(database):
    try:
        tables = request.json['tables']
        connection = create_connection(database)
        cursor = connection.cursor()

        for table_name in tables:
            query = f"DROP TABLE `{table_name}`"
            cursor.execute(query)

        connection.commit()
        close_connection()

        return "Records deleted successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

from datetime import datetime
from flask import request

@app.route('/create-database', methods=['POST'])
def create_database():
    user_name = request.form.get('userName')
    device_info = request.form.get('deviceInfo')
    chip_info = request.form.get('chipInfo')
    macro_info = request.form.get('macroInfo')
    commit_info = request.form.get('commitInfo')
    description = request.form.get('descriptionOfTest')
    date_created = datetime.now().strftime("%Y%m%d%H%M%S")  

    # Validate that all required fields are present
    if not all([user_name, device_info, chip_info, macro_info, commit_info, description]):
        return jsonify({'message': 'All fields are required.'}), 400

    db_name = f"{user_name}_{device_info}_{chip_info}_{macro_info}_{commit_info}_{description}_{date_created}"

    # Attempt to create the database
    if create_db(db_name):
        return jsonify({'message': f"Database '{db_name}' created successfully."}), 200
    else:
        return jsonify({'message': f"Failed to create Database '{db_name}'."}), 500

@app.route('/download_pptx', methods=['POST'])
def download_pptx():
    template_path = '/home/admin2/webapp_2/pptx_template/template.pptx'

    plots = request.json.get('plots', [])  # Retrieve the Base64 encoded images from the POST request

    prs = Presentation(template_path)  # Open the template PowerPoint file as the base for the new presentation
    
    # Retrieve the Base64 encoded images from the POST request
    plots = request.json.get('plots', [])
    
    for plot_data in plots:
        # Decode each Base64 image
        image_data = base64.b64decode(plot_data.split(",")[-1])
        # Open the image for analysis
        image = Image.open(BytesIO(image_data))
        
        # Choose a slide layout (6 is usually a blank slide)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all shapes (including text boxes) from the slide
        for shape in slide.shapes:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Get the image size
        img_width, img_height = image.size
        # Get the slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Calculate the scaling factor to maintain aspect ratio
        ratio = min(slide_width / img_width, slide_height / img_height)
        new_width = int(img_width * ratio)
        new_height = int(img_height * ratio)
        
        # Center the image
        left = int((slide_width - new_width) / 2)
        top = int((slide_height - new_height) / 2)
        
        # Convert the image data back to a BytesIO object
        img_io = BytesIO(image_data)
        # Add the image to the slide
        slide.shapes.add_picture(img_io, left, top, width=new_width, height=new_height)
    
    # Prepare the presentation to be sent in the response
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    # Set up the response with the correct headers
    response = make_response(pptx_io.getvalue())
    response.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response.headers.set('Content-Disposition', 'attachment; filename="Downloaded_Presentation.pptx"')
    
    return response

@app.route('/getDatabases', methods=['GET'])
def get_databases():
    # Create a database connection
    try:
        conn = create_connection()
        cursor = conn.cursor()
        
        # Use your function to get database names
        databases = get_all_databases(cursor)
        
        # Don't forget to close the cursor and connection when done
        cursor.close()
        conn.close()
        
        # Return the list of databases as a JSON response
        return jsonify(databases)

    except mysql.connector.Error as err:
        # In case of any database connection errors, return an error message
        return jsonify({"error": str(err)}), 500

@app.route('/add-item', methods=['POST'])
def create_item():
    item_id = request.form['item_id']
    item_data = request.form['item_data']
    try:
        response = table.put_item(
            Item={
                'ID': item_id,
                'item_data': item_data
            }
        )
        return redirect(url_for('home'))
    except ClientError as e:
        return jsonify({'error': str(e)}), 500

@app.route('/items', methods=['GET'])
def list_items():
    try:
        response = table.scan()
        items = response['Items']
        return jsonify(items)
    except ClientError as e:
        return jsonify({'error': str(e)}), 500

@app.route('/generate_zip', methods=['POST'])
def generate_zip():
    data = request.get_json()
    table_names = data.get('tableNames', [])

    if not table_names:
        return 'No tables selected', 400

    # Create a ZIP file in memory
    memory_file = io.BytesIO()
    with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zf:
        for table_name in table_names:
            # Generate CSV content for each table (replace with your logic)
            df = pd.read_csv(f'data/{table_name}.csv')  # Replace with your actual data source
            csv_content = df.to_csv(index=False)
            zf.writestr(f'{table_name}.csv', csv_content)

    memory_file.seek(0)
    return send_file(memory_file, attachment_filename='tables.zip', as_attachment=True)

def rename_duplicate_columns(df):
    cols = pd.Series(df.columns)
    duplicates = df.columns[df.columns.duplicated()].unique()
    for dup in duplicates:
        cols[df.columns.get_loc(dup)] = [f"{dup}_{i}" if i != 0 else dup for i in range(sum(df.columns == dup))]
    df.columns = cols
    return df

@app.route('/mergeTablesInput', methods=['POST'])
def merge_tables_input():
    database = request.form.get('database')
    table_names = request.form.getlist('tableNames')
    if not database or not table_names:
        return 'Database or table names not provided', 400
    return render_template('input_form_merge.html', database=database, table_names=table_names)

@app.route('/mergeTablesProcess', methods=['POST'])
def merge_tables_process():
    database = request.form.get('database')
    table_names = request.form.getlist('tableNames')
    state_pattern = request.form.get('state_pattern')
    new_table_name = request.form.get('newTableName')

    if not database:
        return 'Database not specified', 400
    if not table_names:
        return 'No tables specified', 400
    if not state_pattern:
        return 'State pattern not specified', 400
    if not new_table_name:
        return 'New table name not specified', 400

    try:
        # Create a connection specifying the database
        connection = create_connection(database)
        cursor = connection.cursor()
        
        # List to hold reshaped arrays
        reshaped_arrays = []

        # Load the pattern file based on state_pattern
        pattern_files = {
            "1296x64_rowbar_4states": "/home/admin2/webapp_2/State_pattern_files/1296x64_rowbar_4states.npy",
            "248x248_checkerboard_4states": "/home/admin2/webapp_2/State_pattern_files/248x248_checkerboard_4states.npy",
            "1296x64_Adrien_random_4states": "/home/admin2/webapp_2/State_pattern_files/1296x64_Adrien_random_4states.npy",
            "248x248_1state": "/home/admin2/webapp_2/State_pattern_files/248x248_1state.npy",
            "1296x64_1state": "/home/admin2/webapp_2/State_pattern_files/1296x64_1state.npy",
            "248x248_16states": "/home/admin2/webapp_2/State_pattern_files/248x248_16states.npy",
            "248x1_1state": "/home/admin2/webapp_2/State_pattern_files/248x1_1state.npy"
        }

        file_path = pattern_files.get(state_pattern)
        if not file_path or not os.path.exists(file_path):
            return 'Invalid state pattern or file not found', 400

        # Load the pattern array and flatten it
        pattern_array = np.load(file_path)
        a, b = pattern_array.shape
        pattern_flat = pattern_array.flatten()

        for table_name in table_names:
            query = f"SELECT * FROM `{table_name}`"
            cursor.execute(query)
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description]
            df = pd.DataFrame(rows, columns=columns)

            # Convert DataFrame to numpy array
            arr = df.to_numpy()

            # Flatten the array
            arr_flat = arr.flatten()

            # Ensure that arr_flat and pattern_flat are the same size
            if arr_flat.size != pattern_flat.size:
                return 'Pattern file and table data dimensions do not match', 400

            # Get indices that would sort the pattern_flat
            pattern_indices = np.argsort(pattern_flat)

            # Reorder arr_flat according to pattern_indices
            arr_reordered = arr_flat[pattern_indices]

            # Reshape to (a*b, 1)
            arr_new = arr_reordered.reshape((a * b, 1))

            # Append to list
            reshaped_arrays.append(arr_new)

        if reshaped_arrays:
            # Concatenate all reshaped arrays along axis=1
            combined_array = np.concatenate(reshaped_arrays, axis=1)

            # Convert back to DataFrame
            combined_df = pd.DataFrame(combined_array)

            # Rename duplicate columns
            combined_df = rename_duplicate_columns(combined_df)

            # Check if the new table name already exists
            cursor.execute("SHOW TABLES LIKE %s", (new_table_name,))
            if cursor.fetchone():
                connection.close()
                return 'A table with the new name already exists.', 400

            # Create the new table in the database using SQLAlchemy engine
            engine = create_db_engine(database)
            combined_df.to_sql(new_table_name, engine, if_exists='fail', index=False)

            # Clean up
            cursor.close()
            close_connection()

            # After successful merging, redirect to list_tables
            return redirect(url_for('list_tables', database=database))

        else:
            # Clean up
            cursor.close()
            close_connection()
            return 'No tables were reshaped and combined.', 400

    except Exception as e:
        return str(e), 500

@app.route('/copy_tables', methods=['POST'])
def copy_tables():
    data = request.get_json()
    source_db = data.get('sourceDatabase')
    target_db = data.get('targetDatabase')
    table_names = data.get('tableNames')

    # Check for missing data
    if not all([source_db, target_db, table_names]):
        return jsonify({'message': 'Missing data in request.'}), 400

    # Connect to databases
    source_conn = create_connection(source_db)
    target_conn = create_connection(target_db)

    try:
        # Get existing table names in the target database
        existing_tables = get_table_names(target_conn)

        # Find conflicts
        conflicts = set(table_names) & set(existing_tables)
        if conflicts:
            conflict_list = ', '.join(conflicts)
            return jsonify({'message': f'The following tables already exist in the target database: {conflict_list}'}), 400

        # Copy tables
        for table in table_names:
            # Fetch table data from the source database
            table_data = get_table_from_database(source_db, table)
            if table_data is None:
                return jsonify({'message': f'Failed to retrieve data for table {table}.'}), 500

            data_rows, columns = table_data

            # Create the table in the target database
            create_table(target_conn, table, data_rows, columns)

        return jsonify({'message': 'Tables copied successfully.'}), 200

    except Exception as e:
        print(f'Error copying tables: {e}')
        return jsonify({'message': 'An error occurred while copying tables.'}), 500

    finally:
        # Close database connections
        source_conn.close()
        target_conn.close()

@app.route('/concatenate_tables', methods=['POST'])
def concatenate_tables():
    data = request.json
    database = data.get('database')
    table_names = data.get('tableNames')
    new_table_name = data.get('newTableName')

    if not (database and table_names and new_table_name):
        return jsonify(success=False, message='Missing required information.')

    try:
        # Connect to the MySQL database using your existing function
        connection = create_connection(database)
        cursor = connection.cursor()

        # Fetch data from each table
        df_list = []
        for table_name in table_names:
            cursor.execute(f"SELECT * FROM `{table_name}`")
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description]
            df = pd.DataFrame(rows, columns=columns)
            df_list.append(df)

        # Verify that all dataframes have the same number of rows
        nrows = df_list[0].shape[0]
        if not all(df.shape[0] == nrows for df in df_list):
            connection.close()
            return jsonify(success=False, message='Selected tables do not have the same number of rows.')

        # Concatenate dataframes column-wise
        concatenated_df = pd.concat(df_list, axis=1)

        # Rename duplicate columns
        concatenated_df = rename_duplicate_columns(concatenated_df)

        # Check if the new table name already exists
        cursor.execute("SHOW TABLES LIKE %s", (new_table_name,))
        if cursor.fetchone():
            connection.close()
            return jsonify(success=False, message='A table with the new name already exists.')

        # Save the concatenated dataframe to a new table using SQLAlchemy engine
        engine = create_db_engine(database)
        concatenated_df.to_sql(new_table_name, con=engine, if_exists='fail', index=False)

        connection.close()

        return jsonify(success=True)
    except Exception as e:
        print(f'Error: {e}')
        return jsonify(success=False, message=str(e))

@app.route('/rename-table', methods=['POST'])
def rename_table():
    data = request.get_json()
    database = data.get('database')
    old_name = data.get('old_name')
    new_name = data.get('new_name')

    if not (database and old_name and new_name):
        return 'Missing required information', 400

    try:
        # Call function to rename table in the database
        result = rename_table_in_database(database, old_name, new_name)
        if result:
            return 'Table renamed successfully', 200
        else:
            return 'Error renaming table', 400
    except Exception as e:
        print(f'Error in rename_table route: {e}')
        return str(e), 500

@app.route('/run-forming-progress')
def run_forming_progress():
    try:
        # Run the script and capture the output
        script_path = '/home/admin2/webapp_2/postprocess/forming_progress.py'
        process = subprocess.Popen(
            [sys.executable, script_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            cwd='/home/admin2/webapp_2/postprocess'  # Optional: set working directory
        )
        stdout, stderr = process.communicate()

        if process.returncode != 0:
            output = f"Script exited with return code {process.returncode}\n"
            output += f"Standard Output:\n{stdout}\n"
            output += f"Standard Error:\n{stderr}\n"
        else:
            output = stdout

        # Render the output in a template
        return render_template('forming_progress_output.html', output=output)

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        return f"An error occurred while running the script: {e}\n{error_details}", 500

@app.route('/simple_combine', methods=['POST'])
def simple_combine():
    data = request.get_json()
    database = data.get('database')
    table_name = data.get('tableName')
    new_table_name = data.get('newTableName')

    if not (database and table_name and new_table_name):
        return jsonify(success=False, message='Missing required information.')

    try:
        # Connect to the database
        connection = create_connection(database)
        cursor = connection.cursor()

        # Fetch data from the table
        cursor.execute(f"SELECT * FROM `{table_name}`")
        rows = cursor.fetchall()
        columns = [desc[0] for desc in cursor.description]
        df = pd.DataFrame(rows, columns=columns)

        # Stack all columns into one Series and reset the index
        combined_series = df.stack().reset_index(drop=True)

        # Create a new DataFrame with the combined data
        combined_df = pd.DataFrame({'Combined': combined_series})

        # Check if the new table name already exists
        cursor.execute("SHOW TABLES LIKE %s", (new_table_name,))
        if cursor.fetchone():
            connection.close()
            return jsonify(success=False, message='A table with the new name already exists.')

        # Save the combined DataFrame to a new table using SQLAlchemy engine
        engine = create_db_engine(database)
        combined_df.to_sql(new_table_name, con=engine, if_exists='fail', index=False)

        connection.close()

        return jsonify(success=True)
    except Exception as e:
        print(f"Error: {e}")
        return jsonify(success=False, message='An error occurred while combining the table.')

@app.route('/check_single_column_tables', methods=['POST'])
def check_single_column_tables():
    data = request.get_json()
    database = data.get('database')
    table_names = data.get('tableNames')

    if not (database and table_names):
        return jsonify(success=False, message='Missing required information.')

    try:
        connection = create_connection(database)
        cursor = connection.cursor()

        for table_name in table_names:
            cursor.execute(f"SELECT COUNT(*) FROM INFORMATION_SCHEMA.COLUMNS WHERE table_schema = %s AND table_name = %s", (database, table_name))
            column_count = cursor.fetchone()[0]
            if column_count != 1:
                connection.close()
                return jsonify(success=False, message=f"Table '{table_name}' does not have exactly one column.")

        connection.close()
        return jsonify(success=True)
    except Exception as e:
        print(f'Error: {e}')
        return jsonify(success=False, message='An error occurred while checking tables.')

@app.route('/combine_single_columns', methods=['POST'])
def combine_single_columns():
    data = request.get_json()
    database = data.get('database')
    table_names = data.get('tableNames')
    new_table_name = data.get('newTableName')

    if not (database and table_names and new_table_name):
        return jsonify(success=False, message='Missing required information.')

    try:
        # Connect to the database
        connection = create_connection(database)
        cursor = connection.cursor()

        combined_data = pd.DataFrame()

        for table_name in table_names:
            # Fetch the single column from the table
            cursor.execute(f"SELECT * FROM `{table_name}`")
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description]
            df = pd.DataFrame(rows, columns=columns)
            column_name = table_name  # Use table name as the column header
            combined_data[column_name] = df.iloc[:, 0]  # Assuming the first column

        # Check if the new table name already exists
        cursor.execute("SHOW TABLES LIKE %s", (new_table_name,))
        if cursor.fetchone():
            connection.close()
            return jsonify(success=False, message='A table with the new name already exists.')

        # Save the combined DataFrame to a new table using SQLAlchemy engine
        engine = create_db_engine(database)
        combined_data.to_sql(new_table_name, con=engine, if_exists='fail', index=False)

        connection.close()
        return jsonify(success=True)
    except Exception as e:
        print(f"Error: {e}")
        return jsonify(success=False, message='An error occurred while combining the tables.')

@app.route('/check_two_column_table', methods=['POST'])
def check_two_column_table():
    data = request.get_json()
    database = data.get('database')
    table_name = data.get('tableName')

    if not (database and table_name):
        return jsonify(success=False, message='Missing required information.')

    try:
        connection = create_connection(database)
        cursor = connection.cursor()

        cursor.execute("""
            SELECT COUNT(*)
            FROM INFORMATION_SCHEMA.COLUMNS
            WHERE table_schema = %s AND table_name = %s
        """, (database, table_name))
        column_count = cursor.fetchone()[0]

        if column_count != 2:
            connection.close()
            return jsonify(success=False, message=f"Table '{table_name}' does not have exactly two columns.")

        connection.close()
        return jsonify(success=True)
    except Exception as e:
        print(f'Error: {e}')
        return jsonify(success=False, message='An error occurred while checking the table.')

@app.route('/generate_scatter_plot/<database>/<table_name>', methods=['GET'])
def generate_scatter_plot(database, table_name):
    try:
        x_column = request.args.get('x_column')
        y_column = request.args.get('y_column')

        if not x_column or not y_column:
            return "X and Y columns must be specified.", 400

        connection = create_connection(database)
        cursor = connection.cursor()

        # Fetch data from the table using the selected columns
        query = f"SELECT `{x_column}`, `{y_column}` FROM `{table_name}`"
        cursor.execute(query)
        rows = cursor.fetchall()
        columns = [desc[0] for desc in cursor.description]

        df = pd.DataFrame(rows, columns=columns)

        # Generate scatter plot
        fig, ax = plt.subplots(figsize=(12, 12))
        ax.scatter(df[x_column], df[y_column],
                   label=f'{y_column} vs {x_column}',
                   color='blue',
                   s=1)  # Adjust 's' to make dots smaller

        ax.set_xlabel(x_column)
        ax.set_ylabel(y_column)
        ax.set_title(f"Scatter Plot of {table_name}")
        ax.legend()

        # Set the aspect ratio of the plot to be equal
        ax.set_aspect('equal', adjustable='box')

        # Convert plot to PNG image
        img = BytesIO()
        plt.savefig(img, format='png')
        img.seek(0)

        plt.close(fig)
        connection.close()

        return send_file(img, mimetype='image/png')
    except Exception as e:
        print(f'Error: {e}')
        return 'An error occurred while generating the scatter plot.', 500

@app.route('/get_table_columns', methods=['POST'])
def get_table_columns():
    data = request.get_json()
    database = data.get('database')
    table_name = data.get('tableName')

    if not (database and table_name):
        return jsonify(success=False, message='Missing required information.')

    try:
        connection = create_connection(database)
        cursor = connection.cursor()
        cursor.execute(f"SHOW COLUMNS FROM `{table_name}`")
        columns = [row[0] for row in cursor.fetchall()]
        connection.close()
        return jsonify(success=True, columns=columns)
    except Exception as e:
        print(f'Error: {e}')
        return jsonify(success=False, message='An error occurred while fetching columns.')

@app.route('/merge-schemas', methods=['POST'])
def merge_schemas():
    try:
        data = request.get_json()
        new_schema_name = data.get('newSchemaName')
        selected_schemas = data.get('selectedSchemas')

        if not new_schema_name or not selected_schemas or len(selected_schemas) < 2:
            return jsonify({'success': False, 'message': 'Invalid input parameters'})

        # Create new schema
        connection = create_connection()
        cursor = connection.cursor()
        
        # Check if new schema already exists
        cursor.execute("SHOW DATABASES LIKE %s", (new_schema_name,))
        if cursor.fetchone():
            return jsonify({'success': False, 'message': 'A folder with this name already exists'})

        # Create new schema
        cursor.execute(f"CREATE DATABASE `{new_schema_name}`")

        # Process each selected schema
        for schema in selected_schemas:
            # Get all tables from current schema
            cursor.execute(f"SHOW TABLES FROM `{schema}`")
            tables = cursor.fetchall()
            
            for (table_name,) in tables:
                # Create new table name with schema prefix
                new_table_name = f"{schema}_{table_name}"
                
                # Copy table structure and data to new schema
                cursor.execute(f"""
                    CREATE TABLE `{new_schema_name}`.`{new_table_name}` 
                    LIKE `{schema}`.`{table_name}`
                """)
                cursor.execute(f"""
                    INSERT INTO `{new_schema_name}`.`{new_table_name}`
                    SELECT * FROM `{schema}`.`{table_name}`
                """)

        connection.commit()
        cursor.close()
        connection.close()

        return jsonify({'success': True, 'message': 'Folders merged successfully'})

    except Exception as e:
        print(f"Error in merge_schemas: {str(e)}")
        return jsonify({'success': False, 'message': str(e)})

#------------------------------------------------------------------------------------------------------------

def sanitize_table_name(name):
    """
    Sanitize the filename to make it suitable for usage as a MySQL table name.
    """
    # Remove non-word characters and spaces
    sanitized_name = re.sub(r'\W+| ', '_', name)

    # Ensure it starts with a letter, prepend an 'a' if not
    #if not sanitized_name[0].isalpha():
        #sanitized_name = 'a' + sanitized_name

    return sanitized_name.lower()

def validate_filename(filename):
    print("Filename:", filename)
    #pattern = r"^lot[A-Za-z0-9]+_wafer[A-Za-z0-9]+_die[A-Za-z0-9]+_dut[A-Za-z0-9]+_[A-Za-z0-9]+_[A-Za-z0-9]+_[A-Za-z0-9]+_[A-Za-z0-9]+\.(csv|txt)$"
    pattern = r"^lot[A-Za-z0-9]+_wafer[A-Za-z0-9]+_die[A-Za-z0-9]+_dut[A-Za-z0-9]*_[A-Za-z0-9]+_[A-Za-z0-9]+_[A-Za-z0-9]+_[A-Za-z0-9]+\.(csv|txt|npy)$"
    print("Pattern match result:", bool(re.match(pattern, filename)))
    return bool(re.match(pattern, filename))

def render_results(results):
    results_html = "<html><body style='background-color: white;'>"
    for result in results:
        results_html += f"<p>{result}</p>"
    results_html += "</body></html>"
    return render_template_string(results_html), 200

def get_form_data_generate_plot(form):
    form_data = {
        key: form.get(key, "").strip() for key in [
            'state_pattern_type', 'number_of_states',
            'selected_groups_1D', 'pass_range_1D', 'state_pattern',
            'selected_groups_predefined', 'pass_range_predefined',
            'custom_selected_groups_predefined', 'custom_pass_range_predefined',
            'color_map_flag', 'outlier_analysis_flag', 'target_values'  # Added target_values here
        ]
    }

    print("Form Data Retrieved:", form_data)  # Debug print

    if form_data['state_pattern_type'] == '1D':
        form_data['number_of_states'] = int(form_data['number_of_states']) if form_data['number_of_states'].isdigit() else None
        # Handle selected_groups
        if form_data.get('selected_groups_1D') == 'custom' and form.get('custom_selected_groups_1D'):
            selected_groups_1D = form.get('custom_selected_groups_1D').split(',')
        else:
            selected_groups_1D = form_data.get('selected_groups_1D', '').split(',') if form_data.get('selected_groups_1D') else []
        form_data['selected_groups'] = [int(float(num)) for num in selected_groups_1D if num.strip()]
    elif form_data['state_pattern_type'] == 'predefined':
        # Handle selected_groups
        if form_data.get('selected_groups_predefined') == 'custom' and form_data.get('custom_selected_groups_predefined'):
            selected_groups_predefined = form_data.get('custom_selected_groups_predefined', '').split(',')
        else:
            selected_groups_predefined = form_data.get('selected_groups_predefined', '').split(',') if form_data.get('selected_groups_predefined') else []
        
        form_data['selected_groups'] = [int(float(num)) for num in selected_groups_predefined if num.strip()]
        print('selected_groups:',  form_data['selected_groups'])

    form_data['state_pattern'] = form.get('state_pattern', None)
    
    # Handle pass_range
    if form_data.get('pass_range_predefined') == 'custom' and form_data.get('custom_pass_range_predefined'):
        pass_range_predefined = form_data.get('custom_pass_range_predefined', '').split(',')
    else:
        pass_range_predefined = form_data.get('pass_range_predefined', '').split(',') if form_data.get('pass_range_predefined') else []
    
    form_data['pass_range'] = [float(num) for num in pass_range_predefined if num.strip()]

    # Convert checkbox flags to boolean
    form_data['color_map_flag'] = form_data.get('color_map_flag', 'False') == 'True'
    form_data['outlier_analysis_flag'] = form_data.get('outlier_analysis_flag', 'False') == 'True'

    # Process target values
    target_values_str = form_data.get('target_values', '')
    if target_values_str:
        try:
            form_data['target_values'] = [float(x.strip()) for x in target_values_str.split(',') if x.strip()]
        except ValueError:
            form_data['target_values'] = []
    else:
        form_data['target_values'] = []

    print("Final Form Data:", form_data)  # Debug print
    return form_data

def get_form_data_generate_plot_ber_by_bls(form):
    return form_data

def get_form_data_generate_plot_read_stability(form):
    # Initialize form_data dictionary
    form_data = {}

    # Retrieve and store the 'state_pattern' from the form
    form_data['state_pattern'] = form.get('state_pattern', None)
    form_data['input_integer'] = form.get('input_integer', None)

    # Debug print to check the retrieved 'state_pattern'
    print("State Pattern:", form_data['state_pattern'])
    print("input_integer:", form_data['input_integer'])

    return form_data

def flatten_sections(array_3d):
    """Flatten 16x16 sections from a 3D array."""
    slices = [array_3d[row:row+16, col:col+16, setting].flatten()
              for setting in range(array_3d.shape[2])
              for row in range(0, 64, 16)
              for col in range(0, 64, 16)]
    return np.concatenate(slices)

def process_mat_file(file_content):
    print('process_mat_file')
    df = None
    mat_data = scipy.io.loadmat(BytesIO(file_content))
    for key in mat_data:
        if not key.startswith('__'):
            data = mat_data[key]
            print("data.shape", data.shape)
            if data.shape == (64, 64, 8) or data.shape == (64, 64, 4) or data.shape == (100, 64, 64):
                df = pd.DataFrame(flattened_data)
            elif data.shape == (1, 64, 64):
                flattened_data = data.squeeze()  # Flattening to (64, 64)
                df = pd.DataFrame(flattened_data)
            elif len(data.shape) == 2:  # Check if the data is 2D
                df = pd.DataFrame(data)
            if df is not None:
                break
    return df

def process_h5py_file(file_stream):  #64x64
    print('process_h5py_file')
    df = None
    with h5py.File(file_stream, 'r') as f:
        for key in f.keys():
            data = f[key]
            if isinstance(data, h5py.Dataset) and data.shape == (64, 64):
                df = pd.DataFrame(data[:]).T
                break
    return df

def is_hdf5_file(file_stream):
    # Check if the file stream is an HDF5 file
    try:
        h5py.File(file_stream, 'r')
        return True
    except OSError:
        return False

import pandas as pd
from pyexcel_xls import get_data
import sys
import json

def process_file(file_stream, file_extension, db_name):
    #print("Python version:", sys.version)
    #print("pandas version:", pd.__version__)
    
    file_stream.seek(0)
    df = None

    if file_extension == "csv":
        try:
            df = pd.read_csv(file_stream, header=None, skiprows=0)
            print("shape of df:", df.shape)
        except Exception as e:
            print("An error occurred:", e)

    elif file_extension == "xlsx":
        print(f"file_extension == {file_extension}") 
        try:
            df = pd.read_excel(file_stream, header=None, skiprows=0, engine='openpyxl')
            print("shape of df:", df.shape)
        except Exception as e:
            print("An error occurred:", e)

    elif file_extension == "xls":
        print(f"file_extension == {file_extension}") 
        try:
            # Read the .xls file using pyexcel
            data = get_data(file_stream)
            # Assuming you want to read the first sheet
            sheet_name = list(data.keys())[0]
            sheet_data = data[sheet_name]
            df = pd.DataFrame(sheet_data)
            print("shape of df:", df.shape)
        except Exception as e:
            print("An error occurred:", e)

    elif file_extension == "txt":
        print("uploading txt")
        content = file_stream.read()
        print(f"Content length: {len(content)}")  # Log the length of the content

        if len(content) > 65535:
            raise ValueError("Content too large to fit in the database column")

        df = pd.DataFrame({'content': [content]})
        print("DataFrame created")  # Log DataFrame creation
        print(df)  # Print the DataFrame

    elif file_extension == "npy":
        try:
            # Load the .npy file into a numpy array
            np_array = np.load(file_stream, allow_pickle=True)
            print("Original array shape:", np_array.shape)

            if np_array.ndim == 1:
                df = pd.DataFrame(np_array)
            elif np_array.ndim == 2:
                df = pd.DataFrame(data=np_array)
                if df.shape[1] > 1017:  #df.shape[1] is the num of columns
                    df = df.transpose()  #to transpose 64x1296 to 1296x64              
            elif np_array.ndim == 3:
                df = pd.DataFrame(data=np_array.reshape(np_array.shape[0], -1))
            elif np_array.ndim == 4:
                squeezed_array = np.squeeze(np_array)
                if squeezed_array.ndim < 4:
                    np_array = squeezed_array
                    print("Array was squeezed to dimensions:", np_array.shape)
                else:
                    print("Squeezing did not reduce dimensions, handling as 4D array.")
                df = pd.DataFrame(data=np_array.reshape(-1, np_array.shape[-1]))
            else:
                raise ValueError("Numpy array dimensionality not supported")

            # Check for NaN values immediately after DataFrame creation
            if df.isnull().values.any():
                print("DataFrame contains NaN values after creation.")
            else:
                print("DataFrame does not contain NaN values after creation.")

            print("DataFrame dtypes:")
            print(df.dtypes)

        except Exception as e:
            print(f"Error processing .npy file: {e}")
            df = pd.DataFrame()

    elif file_extension == "mat":
        # Reset the file stream to the beginning for reading
        file_stream.seek(0)
        if is_hdf5_file(file_stream):
            # If it's an HDF5 file, use the h5py processor
            df = process_h5py_file(file_stream)
        else:
            # For other .mat files, process them here
            try:
                # Reset the file stream again as is_hdf5_file may have moved it
                file_stream.seek(0)
                file_content = file_stream.read()
                df = process_mat_file(file_content)
                if df is None or df.empty:
                    raise ValueError("No suitable dataset found in the .mat file")
            except Exception as e:
                print(f"Error processing .mat file: {e}")
                df = pd.DataFrame()
    
    elif file_extension == "json":
        print(f"file_extension == {file_extension}")
        try:
            file_stream.seek(0)
            df = pd.read_json(file_stream)
            print("shape of df:", df.shape)
        except ValueError as e:
            print("An error occurred with pd.read_json:", e)
            try:
                file_stream.seek(0)
                content = file_stream.read()
                # If content is bytes, decode to string
                if isinstance(content, bytes):
                    content = content.decode('utf-8')
                # Load JSON data
                data = json.loads(content)
                # Normalize JSON data to a flat table
                df = pd.json_normalize(data)
                print("shape of df after normalization:", df.shape)
            except Exception as e:
                print("An error occurred while processing JSON data:", e)
                df = pd.DataFrame()
        except Exception as e:
            print("An unexpected error occurred:", e)
            df = pd.DataFrame()

    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")

    # Fallback for any unprocessed or empty data frames
    if df is None or df.empty:
        print("No data processed for the file, returning empty DataFrame")
        df = pd.DataFrame()
    
    return df

# Add custom Jinja2 filter for NaN values
@app.template_filter('is_nan')
def is_nan_filter(value):
    try:
        return np.isnan(value)
    except:
        return False